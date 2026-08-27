"""Shared slide furniture for the review decks.

Every measurement below was taken from the review-deck template these slides
have to match: 13.333 x 7.5 in (16:9), Times New Roman, centred underlined
headings, a grey Calibri footer and the college logo in the top-left corner.
See template-notes.md for how the values were measured.

A build script calls configure() first, then uses the slide builders.
"""
import os
import copy
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.oxml.ns import qn
from PIL import Image

HERE = os.path.dirname(os.path.abspath(__file__))
LOGO = os.path.join(HERE, "assets", "logo.png")

DIAG = None            # set by configure()

# ── design tokens, measured from the template ───────────────────────────
SLIDE_W, SLIDE_H = 13.333, 7.5

ACCENT_NAVY = RGBColor(0x00, 0x20, 0x60)
TEXT_BLACK = RGBColor(0x00, 0x00, 0x00)
FOOTER_GREY = RGBColor(0x89, 0x89, 0x89)
TABLE_BAND = RGBColor(0xED, 0xED, 0xED)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

HEADING_FONT = "Times New Roman"
BODY_FONT = "Times New Roman"
BULLET_FONT = "Arial"
FOOTER_FONT = "Calibri"
MONO_FONT = "Consolas"

LOGO_L, LOGO_T, LOGO_W, LOGO_H = 0.313, 0.120, 0.707, 1.040
HEAD_T, HEAD_H = 0.42, 0.60
BODY_L, BODY_R = 1.278, 12.428
BODY_T, BODY_B = 1.52, 6.72
IMG_L, IMG_R = 0.72, 12.61
IMG_T, IMG_B = 1.42, 6.70
FOOT_T = 6.95

HEAD_SIZE, BODY_SIZE = 24, 18
FOOT_SIZE = 12

REVIEW_LABEL = ""      # e.g. "First Review-2026-27"; set by configure()
EXPORT_DATE = ""       # e.g. "27-08-2026";            set by configure()
TITLE_LINES = []       # the project title, one entry per line
REVIEW_HEADING = ""    # e.g. "Major Project Stage-1 First Review Presentation"

WARNINGS = []


def configure(review_label, export_date, diagram_dir, review_heading, title_lines):
    """Set everything that differs between one review deck and the next."""
    global REVIEW_LABEL, EXPORT_DATE, DIAG, REVIEW_HEADING, TITLE_LINES
    REVIEW_LABEL = review_label
    EXPORT_DATE = export_date
    DIAG = diagram_dir
    REVIEW_HEADING = review_heading
    TITLE_LINES = list(title_lines)
    WARNINGS.clear()


# ── low-level helpers ───────────────────────────────────────────────────
def textbox(slide, l, t, w, h):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return box, tf


def run(p, text, size, bold=False, color=TEXT_BLACK, font=BODY_FONT,
        italic=False, underline=False):
    r = p.add_run()
    r.text = text
    f = r.font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.underline = underline
    f.color.rgb = color
    f.name = font
    return r


def set_bullet(p, char="•", font=BULLET_FONT, indent_in=0.31):
    """Real PowerPoint bullet with a hanging indent, matching the template."""
    pPr = p._p.get_or_add_pPr()
    pPr.set("marL", str(Emu(Inches(indent_in)).emu))
    pPr.set("indent", str(-Emu(Inches(indent_in)).emu))
    for tag in ("a:buNone", "a:buChar", "a:buAutoNum"):
        for el in pPr.findall(qn(tag)):
            pPr.remove(el)
    bu = pPr.makeelement(qn("a:buFont"), {"typeface": font})
    pPr.append(bu)
    bc = pPr.makeelement(qn("a:buChar"), {"char": char})
    pPr.append(bc)


def no_bullet(p):
    pPr = p._p.get_or_add_pPr()
    for tag in ("a:buChar", "a:buAutoNum"):
        for el in pPr.findall(qn(tag)):
            pPr.remove(el)
    pPr.append(pPr.makeelement(qn("a:buNone"), {}))


def check_overflow(name, text_lines, size, width_in, height_in):
    """Rough fit check: python-pptx cannot shrink text, so warn instead."""
    chars_per_line = max(1, int(width_in * 96 / (size * 0.52)))
    lines = 0
    for ln in text_lines:
        lines += max(1, -(-len(ln) // chars_per_line))
    needed = lines * size * 1.35 / 72.0
    if needed > height_in:
        WARNINGS.append("%s: text needs ~%.2f in but the box is %.2f in"
                        % (name, needed, height_in))


# ── slide furniture ─────────────────────────────────────────────────────
def base(prs, heading=None, number=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])          # blank
    slide.shapes.add_picture(LOGO, Inches(LOGO_L), Inches(LOGO_T),
                             Inches(LOGO_W), Inches(LOGO_H))
    if heading:
        _, tf = textbox(slide, 1.30, HEAD_T, SLIDE_W - 2.60, HEAD_H)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run(p, heading, HEAD_SIZE, bold=True, font=HEADING_FONT, underline=True)
    if number is not None:
        foot(slide, number)
    return slide


def foot(slide, number):
    _, tf = textbox(slide, 1.017, FOOT_T, 2.5, 0.32)
    run(tf.paragraphs[0], EXPORT_DATE, FOOT_SIZE, color=FOOTER_GREY, font=FOOTER_FONT)

    _, tf = textbox(slide, 4.3, FOOT_T, 4.7, 0.32)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run(p, REVIEW_LABEL, FOOT_SIZE, color=FOOTER_GREY, font=FOOTER_FONT)

    _, tf = textbox(slide, 11.0, FOOT_T, 1.317, 0.32)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run(p, str(number), FOOT_SIZE, color=FOOTER_GREY, font=FOOTER_FONT)


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text.strip()


# ── slide builders ──────────────────────────────────────────────────────
def title_slide(prs, n, note):
    slide = base(prs, number=n)

    _, tf = textbox(slide, 1.30, 0.83, SLIDE_W - 2.60, 0.90)
    for i, line in enumerate([REVIEW_HEADING, "on"]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        run(p, line, 18, bold=True, color=ACCENT_NAVY, font=HEADING_FONT)

    _, tf = textbox(slide, 0.90, 1.68, SLIDE_W - 1.80, 1.30)
    for i, line in enumerate(TITLE_LINES):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        run(p, line, 28, bold=True, font=HEADING_FONT)

    _, tf = textbox(slide, 1.65, 3.72, 5.0, 2.4)
    p = tf.paragraphs[0]
    run(p, "Presented by", 18, bold=True)
    for line in ["", "Ms. Muskan Sulathana", "Roll No.  :  ____________"]:
        p = tf.add_paragraph()
        run(p, line, 18)

    _, tf = textbox(slide, 7.10, 3.72, 5.4, 2.4)
    p = tf.paragraphs[0]
    run(p, "Guide", 18, bold=True)
    for line in ["", "Rupesh", "Department of EEE",
                 "BVRIT HYDERABAD College of Engineering", "for Women, Hyderabad"]:
        p = tf.add_paragraph()
        run(p, line, 18)

    notes(slide, note)
    return slide


def content_slide(prs, n, heading, bullets, note, size=BODY_SIZE, sub=None,
                  numbered=False):
    slide = base(prs, heading, n)
    top = BODY_T
    if sub:
        _, tf = textbox(slide, BODY_L, top, BODY_R - BODY_L, 0.45)
        run(tf.paragraphs[0], sub, size, italic=True, color=FOOTER_GREY)
        top += 0.55
    _, tf = textbox(slide, BODY_L, top, BODY_R - BODY_L, BODY_B - top)
    gap = 26 if len(bullets) <= 4 else (20 if len(bullets) == 5 else 13)
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        p.line_spacing = 1.10
        if numbered:
            no_bullet(p)
            pPr = p._p.get_or_add_pPr()
            pPr.set("marL", str(Emu(Inches(0.42)).emu))
            pPr.set("indent", str(-Emu(Inches(0.42)).emu))
            run(p, "[%d]  " % (i + 1), size)
        else:
            set_bullet(p)
        run(p, b, size)
    check_overflow(heading, bullets, size, BODY_R - BODY_L, BODY_B - top)
    notes(slide, note)
    return slide


def two_section_slide(prs, n, h1, b1, h2, b2, note, size=BODY_SIZE):
    """Two headings on one slide, the way the template does Problem + Objective."""
    slide = base(prs, h1, n)
    top = BODY_T - 0.10
    for i, b in enumerate(b1):
        _, tf = textbox(slide, BODY_L, top + i * 0.0, BODY_R - BODY_L, 0.01) if False else (None, None)
    _, tf = textbox(slide, BODY_L, top, BODY_R - BODY_L, 2.55)
    for i, b in enumerate(b1):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(16)
        p.line_spacing = 1.10
        set_bullet(p)
        run(p, b, size)
    check_overflow(h1, b1, size, BODY_R - BODY_L, 2.55)

    mid = top + 2.75
    _, tf = textbox(slide, 1.30, mid, SLIDE_W - 2.60, 0.55)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run(p, h2, HEAD_SIZE, bold=True, font=HEADING_FONT, underline=True)

    body_top = mid + 0.78
    _, tf = textbox(slide, BODY_L, body_top, BODY_R - BODY_L, BODY_B - body_top)
    for i, b in enumerate(b2):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(16)
        p.line_spacing = 1.10
        set_bullet(p)
        run(p, b, size)
    check_overflow(h2, b2, size, BODY_R - BODY_L, BODY_B - body_top)
    notes(slide, note)
    return slide


def _fit(img, l, t, w, h):
    with Image.open(img) as im:
        iw, ih = im.size
    scale = min(w / iw, h / ih)
    dw, dh = iw * scale, ih * scale
    return l + (w - dw) / 2, t + (h - dh) / 2, dw, dh


def image_slide(prs, n, heading, image, note, caption=None, bullets=None,
                bullet_size=17):
    slide = base(prs, heading, n)
    top, bottom = IMG_T, IMG_B
    if bullets:
        _, tf = textbox(slide, BODY_L, IMG_T, BODY_R - BODY_L, 1.0)
        for i, b in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_after = Pt(6)
            set_bullet(p)
            run(p, b, bullet_size)
        top = IMG_T + 0.42 * len(bullets) + 0.22
    if caption:
        bottom -= 0.42
    l, t, w, h = _fit(os.path.join(DIAG, image), IMG_L, top, IMG_R - IMG_L, bottom - top)
    slide.shapes.add_picture(os.path.join(DIAG, image), Inches(l), Inches(t),
                             Inches(w), Inches(h))
    if caption:
        _, tf = textbox(slide, BODY_L, IMG_B - 0.36, BODY_R - BODY_L, 0.34)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run(p, caption, 14, italic=True, color=FOOTER_GREY)
    notes(slide, note)
    return slide


def two_image_slide(prs, n, heading, img1, img2, cap1, cap2, note, bullets=None):
    slide = base(prs, heading, n)
    half = (IMG_R - IMG_L - 0.45) / 2
    img_h = 4.55 - (0.44 * len(bullets) + 0.20 if bullets else 0)
    for i, (img, cap) in enumerate(((img1, cap1), (img2, cap2))):
        left = IMG_L + i * (half + 0.45)
        l, t, w, h = _fit(os.path.join(DIAG, img), left, IMG_T + 0.32, half, img_h)
        slide.shapes.add_picture(os.path.join(DIAG, img), Inches(l), Inches(t),
                                 Inches(w), Inches(h))
        _, tf = textbox(slide, left, IMG_T - 0.06, half, 0.34)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run(p, cap, 16, bold=True)
    if bullets:
        bt = IMG_T + 0.32 + img_h + 0.18
        _, tf = textbox(slide, BODY_L, bt, BODY_R - BODY_L, BODY_B - bt)
        for i, b in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_after = Pt(8)
            set_bullet(p)
            run(p, b, 17)
        check_overflow(heading, bullets, 17, BODY_R - BODY_L, BODY_B - bt)
    notes(slide, note)
    return slide


def table_slide(prs, n, heading, headers, rows, note, widths,
                size=15, total_row=False, sub=None, row_h=0.40, aligns=None):
    slide = base(prs, heading, n)
    top = BODY_T
    if sub:
        _, tf = textbox(slide, BODY_L, top - 0.06, BODY_R - BODY_L, 0.40)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run(p, sub, 14, italic=True, color=FOOTER_GREY)
        top += 0.42

    total_w = sum(widths)
    left = (SLIDE_W - total_w) / 2
    nrows = len(rows) + 1
    height = min(BODY_B - top, row_h * nrows)
    shape = slide.shapes.add_table(nrows, len(headers), Inches(left), Inches(top),
                                   Inches(total_w), Inches(height))
    table = shape.table
    table.first_row = True
    table.horz_banding = False
    for i, w in enumerate(widths):
        table.columns[i].width = Inches(w)

    for c, htext in enumerate(headers):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT_NAVY
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = cell.margin_right = Inches(0.07)
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run(p, htext, size, bold=True, color=WHITE)

    for r, row in enumerate(rows, start=1):
        banded = (r % 2 == 0)
        is_total = total_row and r == len(rows)
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = TABLE_BAND if (banded or is_total) else WHITE
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = cell.margin_right = Inches(0.07)
            cell.margin_top = cell.margin_bottom = Inches(0.03)
            p = cell.text_frame.paragraphs[0]
            if aligns:
                p.alignment = PP_ALIGN.CENTER if aligns[c] == "c" else PP_ALIGN.LEFT
            else:
                p.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
            run(p, val, size, bold=is_total)
    notes(slide, note)
    return slide


def code_slide(prs, n, heading, intro, code, tail, note):
    slide = base(prs, heading, n)
    _, tf = textbox(slide, BODY_L, BODY_T, BODY_R - BODY_L, 0.85)
    for i, b in enumerate(intro):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(6)
        set_bullet(p)
        run(p, b, 17)

    top = BODY_T + 0.40 * len(intro) + 0.30
    line_h = 15 * 1.30 / 72.0                      # 15 pt monospace, 1.3 spacing
    code_h = line_h * len(code) + 0.26
    box = slide.shapes.add_textbox(Inches(BODY_L + 0.55), Inches(top),
                                   Inches(BODY_R - BODY_L - 1.10), Inches(code_h))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
    box.line.color.rgb = RGBColor(0xD8, 0xD8, 0xD8)
    tf = box.text_frame
    tf.word_wrap = False
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = tf.margin_right = Inches(0.16)
    tf.margin_top = tf.margin_bottom = Inches(0.10)
    for i, line in enumerate(code):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        no_bullet(p)
        p.space_after = Pt(0)
        run(p, line, 15, font=MONO_FONT)

    bottom = top + code_h + 0.30
    check_overflow(heading + " (tail)", tail, 17, BODY_R - BODY_L, BODY_B - bottom)
    _, tf = textbox(slide, BODY_L, bottom, BODY_R - BODY_L, BODY_B - bottom)
    for i, b in enumerate(tail):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(8)
        set_bullet(p)
        run(p, b, 17)
    notes(slide, note)
    return slide


def closing_slide(prs, n):
    slide = base(prs, number=n)
    _, tf = textbox(slide, 1.30, 2.85, SLIDE_W - 2.60, 1.6)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run(p, "Thank You", 40, bold=True, font=HEADING_FONT)
    p = tf.add_paragraph()
    p.alignment = PP_ALIGN.CENTER
    run(p, "", 18)
    p = tf.add_paragraph()
    p.alignment = PP_ALIGN.CENTER
    run(p, "Questions are welcome", 20, color=ACCENT_NAVY)
    notes(slide, """
That is the end of my presentation. Thank you for listening. I am happy to take
questions on the circuit, the choice of threshold, or the schedule.
""")
    return slide
