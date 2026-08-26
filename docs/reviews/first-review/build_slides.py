"""Build the First Review presentation.

Every measurement below comes from the zeroth-review template that the deck
has to match: 13.333 x 7.5 in (16:9), Times New Roman, centred underlined
headings, a grey Calibri footer and the college logo in the top-left corner.
See ../template-notes.md for how those values were measured.

    python build_slides.py
"""
import os
import copy
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.oxml.ns import qn
from PIL import Image

HERE = os.path.dirname(os.path.abspath(__file__))
DIAG = os.path.join(HERE, "diagrams")
LOGO = os.path.join(HERE, "..", "assets", "logo.png")
OUT = os.path.join(HERE, "Predictive_BMS_First_Review.pptx")

# ── design tokens, measured from the template ───────────────────────────
SLIDE_W, SLIDE_H = 13.333, 7.5

ACCENT_NAVY = RGBColor(0x00, 0x20, 0x60)
TEXT_BLACK = RGBColor(0x00, 0x00, 0x00)
FOOTER_GREY = RGBColor(0x89, 0x89, 0x89)
TABLE_BAND = RGBColor(0xED, 0xED, 0xED)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

HEADING_FONT = "Times New Roman"
BODY_FONT = "Times New Roman"
BULLET_FONT = "Arial"
FOOTER_FONT = "Calibri"
MONO_FONT = "Consolas"

LOGO_L, LOGO_T, LOGO_W, LOGO_H = 0.313, 0.120, 0.707, 1.040
HEAD_T, HEAD_H = 0.42, 0.60
BODY_L, BODY_R = 1.278, 12.428
BODY_T, BODY_B = 1.52, 6.72
IMG_L, IMG_R = 0.72, 12.61
IMG_T, IMG_B = 1.42, 6.70
FOOT_T = 6.95

HEAD_SIZE, BODY_SIZE = 24, 18
FOOT_SIZE = 12

REVIEW_LABEL = "First Review-2026-27"
EXPORT_DATE = "27-08-2026"

WARNINGS = []


# ── low-level helpers ───────────────────────────────────────────────────
def textbox(slide, l, t, w, h):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return box, tf


def run(p, text, size, bold=False, color=TEXT_BLACK, font=BODY_FONT,
        italic=False, underline=False):
    r = p.add_run()
    r.text = text
    f = r.font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.underline = underline
    f.color.rgb = color
    f.name = font
    return r


def set_bullet(p, char="•", font=BULLET_FONT, indent_in=0.31):
    """Real PowerPoint bullet with a hanging indent, matching the template."""
    pPr = p._p.get_or_add_pPr()
    pPr.set("marL", str(Emu(Inches(indent_in)).emu))
    pPr.set("indent", str(-Emu(Inches(indent_in)).emu))
    for tag in ("a:buNone", "a:buChar", "a:buAutoNum"):
        for el in pPr.findall(qn(tag)):
            pPr.remove(el)
    bu = pPr.makeelement(qn("a:buFont"), {"typeface": font})
    pPr.append(bu)
    bc = pPr.makeelement(qn("a:buChar"), {"char": char})
    pPr.append(bc)


def no_bullet(p):
    pPr = p._p.get_or_add_pPr()
    for tag in ("a:buChar", "a:buAutoNum"):
        for el in pPr.findall(qn(tag)):
            pPr.remove(el)
    pPr.append(pPr.makeelement(qn("a:buNone"), {}))


def check_overflow(name, text_lines, size, width_in, height_in):
    """Rough fit check: python-pptx cannot shrink text, so warn instead."""
    chars_per_line = max(1, int(width_in * 96 / (size * 0.52)))
    lines = 0
    for ln in text_lines:
        lines += max(1, -(-len(ln) // chars_per_line))
    needed = lines * size * 1.35 / 72.0
    if needed > height_in:
        WARNINGS.append("%s: text needs ~%.2f in but the box is %.2f in"
                        % (name, needed, height_in))


# ── slide furniture ─────────────────────────────────────────────────────
def base(prs, heading=None, number=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])          # blank
    slide.shapes.add_picture(LOGO, Inches(LOGO_L), Inches(LOGO_T),
                             Inches(LOGO_W), Inches(LOGO_H))
    if heading:
        _, tf = textbox(slide, 1.30, HEAD_T, SLIDE_W - 2.60, HEAD_H)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run(p, heading, HEAD_SIZE, bold=True, font=HEADING_FONT, underline=True)
    if number is not None:
        foot(slide, number)
    return slide


def foot(slide, number):
    _, tf = textbox(slide, 1.017, FOOT_T, 2.5, 0.32)
    run(tf.paragraphs[0], EXPORT_DATE, FOOT_SIZE, color=FOOTER_GREY, font=FOOTER_FONT)

    _, tf = textbox(slide, 4.3, FOOT_T, 4.7, 0.32)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run(p, REVIEW_LABEL, FOOT_SIZE, color=FOOTER_GREY, font=FOOTER_FONT)

    _, tf = textbox(slide, 11.0, FOOT_T, 1.317, 0.32)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run(p, str(number), FOOT_SIZE, color=FOOTER_GREY, font=FOOTER_FONT)


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text.strip()


# ── slide builders ──────────────────────────────────────────────────────
def title_slide(prs, n):
    slide = base(prs, number=n)

    _, tf = textbox(slide, 1.30, 0.83, SLIDE_W - 2.60, 0.90)
    for i, line in enumerate(["Major Project Stage-1 First Review Presentation", "on"]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        run(p, line, 18, bold=True, color=ACCENT_NAVY, font=HEADING_FONT)

    _, tf = textbox(slide, 0.90, 1.68, SLIDE_W - 1.80, 1.30)
    for i, line in enumerate(["Predictive Thermal Battery Management System",
                              "for Li-ion Battery Packs"]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        run(p, line, 28, bold=True, font=HEADING_FONT)

    _, tf = textbox(slide, 1.65, 3.72, 5.0, 2.4)
    p = tf.paragraphs[0]
    run(p, "Presented by", 18, bold=True)
    for line in ["", "Ms. Muskan Sulathana", "Roll No.  :  ____________"]:
        p = tf.add_paragraph()
        run(p, line, 18)

    _, tf = textbox(slide, 7.10, 3.72, 5.4, 2.4)
    p = tf.paragraphs[0]
    run(p, "Guide", 18, bold=True)
    for line in ["", "Rupesh", "Department of EEE",
                 "BVRIT HYDERABAD College of Engineering", "for Women, Hyderabad"]:
        p = tf.add_paragraph()
        run(p, line, 18)

    notes(slide, """
Good morning. I am Muskan Sulathana and this is my first review for the major
project. My project is a battery management system for a small lithium-ion
pack. What makes it different is that it watches how fast the cells are heating
up, not just how hot they are. I will take you through the problem, the circuit
I have designed, the firmware plan, and where I am on the schedule.
""")
    return slide


def content_slide(prs, n, heading, bullets, note, size=BODY_SIZE, sub=None,
                  numbered=False):
    slide = base(prs, heading, n)
    top = BODY_T
    if sub:
        _, tf = textbox(slide, BODY_L, top, BODY_R - BODY_L, 0.45)
        run(tf.paragraphs[0], sub, size, italic=True, color=FOOTER_GREY)
        top += 0.55
    _, tf = textbox(slide, BODY_L, top, BODY_R - BODY_L, BODY_B - top)
    gap = 26 if len(bullets) <= 4 else (20 if len(bullets) == 5 else 13)
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        p.line_spacing = 1.10
        if numbered:
            no_bullet(p)
            pPr = p._p.get_or_add_pPr()
            pPr.set("marL", str(Emu(Inches(0.42)).emu))
            pPr.set("indent", str(-Emu(Inches(0.42)).emu))
            run(p, "[%d]  " % (i + 1), size)
        else:
            set_bullet(p)
        run(p, b, size)
    check_overflow(heading, bullets, size, BODY_R - BODY_L, BODY_B - top)
    notes(slide, note)
    return slide


def _fit(img, l, t, w, h):
    with Image.open(img) as im:
        iw, ih = im.size
    scale = min(w / iw, h / ih)
    dw, dh = iw * scale, ih * scale
    return l + (w - dw) / 2, t + (h - dh) / 2, dw, dh


def image_slide(prs, n, heading, image, note, caption=None, bullets=None,
                bullet_size=17):
    slide = base(prs, heading, n)
    top, bottom = IMG_T, IMG_B
    if bullets:
        _, tf = textbox(slide, BODY_L, IMG_T, BODY_R - BODY_L, 1.0)
        for i, b in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_after = Pt(6)
            set_bullet(p)
            run(p, b, bullet_size)
        top = IMG_T + 0.42 * len(bullets) + 0.22
    if caption:
        bottom -= 0.42
    l, t, w, h = _fit(os.path.join(DIAG, image), IMG_L, top, IMG_R - IMG_L, bottom - top)
    slide.shapes.add_picture(os.path.join(DIAG, image), Inches(l), Inches(t),
                             Inches(w), Inches(h))
    if caption:
        _, tf = textbox(slide, BODY_L, IMG_B - 0.36, BODY_R - BODY_L, 0.34)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run(p, caption, 14, italic=True, color=FOOTER_GREY)
    notes(slide, note)
    return slide


def two_image_slide(prs, n, heading, img1, img2, cap1, cap2, note, bullets=None):
    slide = base(prs, heading, n)
    half = (IMG_R - IMG_L - 0.45) / 2
    img_h = 4.55 - (0.44 * len(bullets) + 0.20 if bullets else 0)
    for i, (img, cap) in enumerate(((img1, cap1), (img2, cap2))):
        left = IMG_L + i * (half + 0.45)
        l, t, w, h = _fit(os.path.join(DIAG, img), left, IMG_T + 0.32, half, img_h)
        slide.shapes.add_picture(os.path.join(DIAG, img), Inches(l), Inches(t),
                                 Inches(w), Inches(h))
        _, tf = textbox(slide, left, IMG_T - 0.06, half, 0.34)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run(p, cap, 16, bold=True)
    if bullets:
        bt = IMG_T + 0.32 + img_h + 0.18
        _, tf = textbox(slide, BODY_L, bt, BODY_R - BODY_L, BODY_B - bt)
        for i, b in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_after = Pt(8)
            set_bullet(p)
            run(p, b, 17)
        check_overflow(heading, bullets, 17, BODY_R - BODY_L, BODY_B - bt)
    notes(slide, note)
    return slide


def table_slide(prs, n, heading, headers, rows, note, widths,
                size=15, total_row=False, sub=None, row_h=0.40, aligns=None):
    slide = base(prs, heading, n)
    top = BODY_T
    if sub:
        _, tf = textbox(slide, BODY_L, top - 0.06, BODY_R - BODY_L, 0.40)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run(p, sub, 14, italic=True, color=FOOTER_GREY)
        top += 0.42

    total_w = sum(widths)
    left = (SLIDE_W - total_w) / 2
    nrows = len(rows) + 1
    height = min(BODY_B - top, row_h * nrows)
    shape = slide.shapes.add_table(nrows, len(headers), Inches(left), Inches(top),
                                   Inches(total_w), Inches(height))
    table = shape.table
    table.first_row = True
    table.horz_banding = False
    for i, w in enumerate(widths):
        table.columns[i].width = Inches(w)

    for c, htext in enumerate(headers):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT_NAVY
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = cell.margin_right = Inches(0.07)
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run(p, htext, size, bold=True, color=WHITE)

    for r, row in enumerate(rows, start=1):
        banded = (r % 2 == 0)
        is_total = total_row and r == len(rows)
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = TABLE_BAND if (banded or is_total) else WHITE
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = cell.margin_right = Inches(0.07)
            cell.margin_top = cell.margin_bottom = Inches(0.03)
            p = cell.text_frame.paragraphs[0]
            if aligns:
                p.alignment = PP_ALIGN.CENTER if aligns[c] == "c" else PP_ALIGN.LEFT
            else:
                p.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
            run(p, val, size, bold=is_total)
    notes(slide, note)
    return slide


def code_slide(prs, n, heading, intro, code, tail, note):
    slide = base(prs, heading, n)
    _, tf = textbox(slide, BODY_L, BODY_T, BODY_R - BODY_L, 0.85)
    for i, b in enumerate(intro):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(6)
        set_bullet(p)
        run(p, b, 17)

    top = BODY_T + 0.40 * len(intro) + 0.30
    line_h = 15 * 1.30 / 72.0                      # 15 pt monospace, 1.3 spacing
    code_h = line_h * len(code) + 0.26
    box = slide.shapes.add_textbox(Inches(BODY_L + 0.55), Inches(top),
                                   Inches(BODY_R - BODY_L - 1.10), Inches(code_h))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
    box.line.color.rgb = RGBColor(0xD8, 0xD8, 0xD8)
    tf = box.text_frame
    tf.word_wrap = False
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = tf.margin_right = Inches(0.16)
    tf.margin_top = tf.margin_bottom = Inches(0.10)
    for i, line in enumerate(code):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        no_bullet(p)
        p.space_after = Pt(0)
        run(p, line, 15, font=MONO_FONT)

    bottom = top + code_h + 0.30
    check_overflow(heading + " (tail)", tail, 17, BODY_R - BODY_L, BODY_B - bottom)
    _, tf = textbox(slide, BODY_L, bottom, BODY_R - BODY_L, BODY_B - bottom)
    for i, b in enumerate(tail):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(8)
        set_bullet(p)
        run(p, b, 17)
    notes(slide, note)
    return slide


def closing_slide(prs, n):
    slide = base(prs, number=n)
    _, tf = textbox(slide, 1.30, 2.85, SLIDE_W - 2.60, 1.6)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run(p, "Thank You", 40, bold=True, font=HEADING_FONT)
    p = tf.add_paragraph()
    p.alignment = PP_ALIGN.CENTER
    run(p, "", 18)
    p = tf.add_paragraph()
    p.alignment = PP_ALIGN.CENTER
    run(p, "Questions are welcome", 20, color=ACCENT_NAVY)
    notes(slide, """
That is the end of my presentation. Thank you for listening. I am happy to take
questions on the circuit, the choice of threshold, or the schedule.
""")
    return slide


# ── the deck ────────────────────────────────────────────────────────────
def build():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    n = 1

    title_slide(prs, n); n += 1

    content_slide(prs, n, "Contents", [
        "Recap of the zeroth review",
        "Problem statement and objectives",
        "Literature survey",
        "Proposed system, hardware and circuits",
        "Firmware and the predictive algorithm",
        "IoT dashboard, cost and work plan",
    ], """
This is the order I will follow. I will start with what was approved in the
zeroth review, then explain the problem I am solving and why the usual approach
falls short. After that I will go through the circuit design in detail, then the
firmware and the algorithm that does the prediction. I will finish with the
cost, what I have finished so far, and the plan for the rest of the semester.
"""); n += 1

    content_slide(prs, n, "Recap of the Zeroth Review", [
        "The idea was approved: a BMS that acts on how fast a cell heats, not only on how hot it is",
        "The synopsis and the block diagram were accepted",
        "My guide signed off, so Phase 1 of the plan is complete",
        "Since then I have fixed the circuit values and priced every part",
        "This review covers the detailed design and the firmware plan",
    ], """
In the zeroth review I presented the idea and a block diagram, and the panel
approved it. My guide signed off after that, so the approvals phase is finished.
Since then I have done the detailed work: I picked real component values,
worked out what the ADC can actually resolve, and priced the full parts list.
Today I want to show that detail and get feedback before I start building.
"""); n += 1

    content_slide(prs, n, "Problem Statement", [
        "A normal small BMS cuts power only when a fixed temperature is crossed, often 60 °C",
        "A cell that is climbing fast at 35 °C is already faulty, but it passes that test",
        "Once a lithium cell starts heating itself, the heat drives more reaction, which makes more heat",
        "A sensor on the outside of a cell always lags what is happening inside it",
        "So a fixed limit reacts late, and it gives no warning at all beforehand",
    ], """
This is the gap I am working on. Almost every small battery protection board
uses one rule: if the temperature goes above a set value, disconnect. The
problem is that the set value tells you nothing until it is reached. A cell that
is heating one degree every second is clearly in trouble even when it is only at
thirty-five degrees, but a fixed limit ignores that completely. On top of that,
the thermistor sits on the outside of the can, so the inside is always hotter
than what I measure. Waiting for an absolute number means acting late.
"""); n += 1

    content_slide(prs, n, "Objectives", [
        "Build a working 4S battery management system for 18650 cells",
        "Measure four cell voltages, the pack current and four cell temperatures",
        "Trip on the rate at which temperature rises, not only on its value",
        "Switch charging and discharging separately, so only the unsafe direction stops",
        "Send every reading over Wi-Fi to a page that opens on a phone or a laptop",
    ], """
These are the five things I want to have working by the end of the project. The
third one is the new part and the reason for the project. The fourth matters
because a pack that is too cold to charge is still perfectly safe to discharge,
so cutting both directions together would be crude. The last one is what makes
it possible to actually see what the board is doing while it runs.
"""); n += 1

    table_slide(prs, n, "Literature Survey",
        ["Ref", "Author and year", "What they did", "What it does not do", "What I do differently"],
        [["[1]", "Feng et al., 2018",
          "Laboratory study of how runaway starts and spreads",
          "Explains the mechanism; it is not a live detector",
          "Turns that mechanism into a trip that runs on the board"],
         ["[2]", "Zhang et al., 2023",
          "Review of ways to predict runaway early",
          "Most methods need gas or pressure sensors, or a server",
          "One threshold test on the ESP32, with no extra sensors"],
         ["[3]", "Chen et al., 2024",
          "Trips when the temperature rise reaches 1 °C per second",
          "Needs training data and a fitted thermal model",
          "Uses the same 1 °C/s figure, computed directly in firmware"],
         ["[4]", "Habib et al., 2023",
          "Review of BMS functions and the problems still open",
          "Lists thermal runaway as unsolved; builds nothing",
          "Builds and bench-tests one specific answer to it"]],
        """
These four papers frame the work. Feng and his co-authors explain what actually
happens inside a cell as it runs away, and that is where the physics comes from.
Zhang's review shows that most early-warning work needs extra hardware or a
server behind it, which a student project cannot use. The Chen paper is the
closest to what I am doing, and it is where my one degree per second threshold
comes from. Habib's review lists thermal runaway as an open problem but does not
build anything, and that is the gap I am filling.
""",
        widths=[0.62, 1.90, 3.05, 3.05, 3.28], size=13.5, row_h=0.86,
        aligns=["c", "c", "l", "l", "l"]); n += 1

    image_slide(prs, n, "Proposed System", "01_system_block.png", """
This is the whole system on one page. The pack feeds the sensing front end,
which reports voltage, current and temperature to the ESP32. The ESP32 runs the
maths and drives two MOSFETs that sit in the pack's negative return. The same
readings go out over Wi-Fi to a dashboard. The one thing worth pointing out is
the note at the bottom: the ESP32's second ADC block cannot be used while Wi-Fi
is running, so every analogue input has to go to the first block.
""",
        bullets=[
            "The pack, the sensing, the controller and the protection sit in one loop",
            "The same readings are used to protect the pack and to feed the dashboard",
        ]); n += 1

    image_slide(prs, n, "Hardware Architecture", "02_hardware_architecture.png", """
This is the same system with the real parts in it. Four cell taps go through
dividers into the first ADC block. Four thermistors go through a CD4051 analogue
multiplexer, because a DevKit board only exposes six usable ADC1 pins and I need
eight channels. Current is measured by an INA226 over I2C, which costs no ADC
pin at all. A buck module makes the 3.3 volt rail, not a linear regulator,
because dropping thirteen volts in a linear part would waste about two watts.
"""); n += 1

    image_slide(prs, n, "Power Path — Single Line Diagram",
                "03_power_path_sld.png", """
This is the power path on its own. The fuse sits in the positive line as the
last-resort protection if the electronics fail completely. Both MOSFETs sit in
the negative return, which is the same arrangement commercial protection boards
use, because it lets both gates be driven against system ground with no charge
pump. The circles are the measurement points: cell voltages at the taps, pack
current across the shunt, and temperature at the cells themselves.
""",
        bullets=[
            "The 15 A fuse is the backstop if the electronics fail completely",
            "Both switches sit in the negative return, so both gates drive against ground",
        ]); n += 1

    two_image_slide(prs, n, "Sensing Circuits",
        "04_voltage_sense_schematic.png", "05_ntc_schematic.png",
        "Cell voltage", "Cell temperature",
        bullets=[
            "Divider ratio 0.128 puts a full 16.8 V pack at 2.15 V — about 4.7 mV per count",
            "Beta equation: 1/T = 1/T₀ + (1/β)·ln(R/R₀), with T₀ = 298.15 K, R₀ = 10 kΩ, β = 3950 K",
        ], note="""
On the left is the divider for a cell tap. All four taps use the same ratio, so
the top tap sets it: the full pack has to land inside the ADC range, which gives
0.128 and about 4.7 millivolts per count once it is referred back to the tap.
That is not precise enough for a production BMS, which is why I have noted an
external 16-bit converter as the upgrade path. On the right is the thermistor
input. The useful thing here is that a fixed calibration error cancels out when
you take a derivative, so the rate trip is far less sensitive to ADC error than
an absolute reading would be.
"""); n += 1

    image_slide(prs, n, "Protection and Balancing",
                "06_cutoff_balancing_schematic.png", """
The two cut-off MOSFETs are wired source to source. That matters because a
single MOSFET has a body diode which would keep conducting in one direction even
when the device is off. With two facing opposite ways, each one blocks a
different direction, so I can stop charging and discharging separately. I chose
the IRLZ44N because it is a logic-level part, so a five volt gate is enough, and
its on-resistance is about 22 milliohms, which at ten amps is only a couple of
watts. On the right, each cell has its own bleed resistor and switch, and an
optocoupler keeps the drive isolated so one ground-referenced pin can switch a
cell sitting several volts up.
""",
        bullets=[
            "IRLZ44N: logic level, 55 V, about 22 mΩ on-resistance at a 5 V gate",
            "The two body diodes face opposite ways, so each direction blocks separately",
        ]); n += 1

    image_slide(prs, n, "Firmware Architecture", "07_firmware_flowchart.png", """
The firmware runs three FreeRTOS tasks. The sensor task reads everything at ten
hertz and filters it. The safety task owns the state machine and is the only
thing allowed to touch the MOSFET pins, so there is no way for the telemetry
code to accidentally switch the pack. The telemetry task builds the JSON message
once a second. On the right is the state machine. The important part is that
leaving the cut-off state needs both a low rate and a low temperature, held for
five seconds, and after three trips in ten minutes it latches until someone
resets it by hand.
"""); n += 1

    code_slide(prs, n, "The Predictive Algorithm",
        ["Temperature is sampled every 100 ms and smoothed before the slope is taken",
         "The slope is an exponential moving average over a one-second window"],
        ["T_filt  = ema(T_raw, alpha = 0.2)",
         "slope   = (T_filt - T_filt_1s_ago) / 1.0",
         "dTdt    = ema(slope, alpha = 0.3)",
         "",
         "if dTdt >= 1.0 or T_filt >= 60.0:",
         "    open_both_mosfets()",
         "    state = CUTOFF",
         "elif dTdt >= 0.5:",
         "    state = WARNING",
         "# leave CUTOFF only when dTdt < 0.2 and T_filt < 40 for 5 s"],
        ["Smoothing comes first: a slope taken from raw counts is mostly noise",
         "A fixed calibration error cancels in a slope, so the trip tolerates ADC drift"],
        """
This is the core of the project in ten lines. The order matters. If you take the
slope of raw readings you mostly measure noise, because differentiating makes
noise worse, so the smoothing has to come first. Then the slope itself gets
smoothed again over a one second window. The threshold of one degree per second
comes from the Chen paper in my literature survey. The sixty degree line is
still there as a backstop, but in a real runaway the rate trip fires well before
it. The last point is the one I like most: because a derivative removes any
constant offset, a calibration error that would ruin an absolute reading has
almost no effect on the rate.
"""); n += 1

    image_slide(prs, n, "Predictive Compared With a Fixed Limit",
                "08_predictive_vs_reactive.png", """
This is a simulated curve, not measured data, and the slide says so. The cell
warms slowly, then self-heating takes over and the temperature climbs
exponentially. The green line is where the rate trip fires: thirty point eight
seconds, with the cell still at only thirty-six degrees. The red line is where a
fixed sixty degree limit would fire, at forty point four seconds. That is almost
ten seconds earlier, and more importantly the pack is twenty-four degrees cooler
when the power is removed. I will replace this with real bench data once the
prototype is running.
""",
        caption="Simulated for illustration. Real measurements will replace this after bench testing."); n += 1

    image_slide(prs, n, "IoT Dashboard", "09_iot_architecture.png", """
The ESP32 joins the local Wi-Fi and serves a small web page itself, so there is
no cloud account and no broker to set up. Once a second the telemetry task
builds the JSON message you can see at the bottom left, and pushes it over a
WebSocket. The browser draws live gauges and a rolling chart from that. The one
message going the other way is the command to clear a latched cut-off, and that
is deliberately the only thing the dashboard is allowed to do.
"""); n += 1

    table_slide(prs, n, "Bill of Materials",
        ["Item", "Qty", "Unit (₹)", "Amount (₹)"],
        [["ESP32-WROOM-32 DevKit V1", "1", "450", "450"],
         ["18650 cells, 2600 mAh", "4", "250", "1000"],
         ["4S holder and nickel strip", "1 set", "270", "270"],
         ["NTC thermistor 10 kΩ B3950", "4", "7", "28"],
         ["INA226 module with shunt", "1", "250", "250"],
         ["IRLZ44N cut-off MOSFETs", "2", "45", "90"],
         ["Balancing parts per cell (P-FET, opto, bleed)", "4 sets", "45", "180"],
         ["CD4051B multiplexer and TC4420 driver", "1 set", "145", "145"],
         ["MP1584EN buck module", "1", "90", "90"],
         ["Passives, fuse, board and wiring", "1 set", "550", "550"],
         ["Total", "", "", "3053"]],
        """
This is the full parts list. It comes to about three thousand rupees, which is
within what I can fund myself. The thermistor price is the only one I have
confirmed on a live product page; the rest are the usual retail figures and I
have marked them as approximate. The cells are the biggest single line, and I
plan to buy them from a seller who will supply matched capacities, because
mismatched cells would make the balancing work much harder later.
""",
        widths=[5.20, 1.10, 1.60, 1.90], size=14, total_row=True,
        aligns=["l", "c", "c", "c"],
        sub="Approximate Indian retail prices, August 2026"); n += 1

    table_slide(prs, n, "Work Completed",
        ["Phase", "Status", "What is done"],
        [["1 · Initial approvals", "Complete", "Synopsis, block diagram and guide approval"],
         ["2 · Procurement", "In progress", "Parts list finalised and priced; ordering next"],
         ["3 · Prototyping", "Not started", "Waiting on parts"],
         ["4 · Core software", "Not started", "Algorithm and state machine designed on paper"],
         ["5 · IoT integration", "Not started", "Message format decided"],
         ["6 · Final polish", "Not started", "—"]],
        """
This is where I actually am. The approvals phase is finished. Procurement is the
live one: the parts list is final and priced, and I am placing the order this
week. Nothing on the bench has been built yet, which is honest, but the design
work for the later phases is not zero either. The algorithm and the state
machine are worked out on paper, and the telemetry message format is decided, so
those phases should move quickly once the hardware exists.
""",
        widths=[3.10, 2.05, 6.00], size=15, row_h=0.72,
        aligns=["l", "c", "l"]); n += 1

    image_slide(prs, n, "Work Plan", "10_gantt.png", """
This is the plan for the rest of the semester, twelve weeks from the end of
August. Procurement takes two weeks. Prototyping overlaps with it slightly
because I can start the sensing board before the cells arrive. The software
phase is the longest at four and a half weeks, and it deliberately overlaps
prototyping, since I can test the reading and filtering code on the bench supply
before the pack is finished. The four dashed lines are the checkpoints I am
holding myself to.
"""); n += 1

    content_slide(prs, n, "Expected Outcomes", [
        "A working 4S BMS that cuts off on temperature rate, demonstrated on the bench",
        "Measured proof that the rate trip fires earlier than a fixed 60 °C limit",
        "A live dashboard showing voltage, current, temperature and rate for every cell",
        "A design that costs about ₹3,000 and uses only parts available in India",
    ], """
These are the four things I expect to be able to show at the end. The second one
is the real test of the project, and I plan to prove it by warming a cell with a
small heater and recording when each trip fires. The last point matters for a
college project: everything on the list can be bought locally, so the work can
actually be repeated by someone else in the department.
"""); n += 1

    content_slide(prs, n, "References", [
        "X. Feng, M. Ouyang, X. Liu, L. Lu, Y. Xia and X. He, “Thermal runaway mechanism of "
        "lithium ion battery for electric vehicles: A review,” Energy Storage Materials, "
        "vol. 10, pp. 246–267, 2018.",
        "X. Zhang, S. Chen, J. Zhu et al., “A critical review of thermal runaway prediction and "
        "early-warning methods for lithium-ion batteries,” Energy Material Advances, "
        "vol. 4, art. 0008, 2023.",
        "Q. Chen, Y. He, N. Fang and G. Yu, “A combined data-driven and model-based algorithm for "
        "accurate battery thermal runaway warning,” Sensors, vol. 24, no. 15, art. 4964, 2024.",
        "A. K. M. A. Habib, M. K. Hasan, G. F. Issa, D. Singh, S. Islam and T. M. Ghazal, "
        "“Lithium-ion battery management system for electric vehicles: Constraints, challenges "
        "and recommendations,” Batteries, vol. 9, no. 3, art. 152, 2023.",
        "IS 16046 (Part 2) : 2018 / IEC 62133-2 : 2017, Secondary cells and batteries containing "
        "alkaline or other non-acid electrolytes — Part 2: Lithium systems.",
    ], """
These are the five sources I have used. The first four are the papers in my
literature survey table, numbered in the same order. The last one is the Indian
standard that covers sealed lithium cells, and it is the standard whose abuse
tests my design is aimed at, although I am not claiming any certification.
""", size=14, numbered=True); n += 1

    closing_slide(prs, n); n += 1

    prs.save(OUT)
    print("saved %s  (%d slides)" % (os.path.basename(OUT), len(prs.slides.__iter__.__self__._sldIdLst)))
    if WARNINGS:
        print("\noverflow warnings:")
        for w in WARNINGS:
            print("   !!", w)
    else:
        print("no overflow warnings")


if __name__ == "__main__":
    build()
